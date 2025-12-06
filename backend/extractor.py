import pytesseract
from PIL import Image
import io
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation

# Ensure pytesseract can find the tesseract executable if it's not in PATH
# pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

def extract_text_from_image(image_data: bytes, lang='tam+eng') -> str:
    try:
        image = Image.open(io.BytesIO(image_data))
        text = pytesseract.image_to_string(image, lang=lang)
        return text
    except Exception as e:
        print(f"Error extracting text from image: {e}")
        return ""

def extract_from_pdf(file_bytes: bytes) -> str:
    text_content = []
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        for page_num, page in enumerate(doc):
            # Extract text
            text = page.get_text()
            if text:
                text_content.append(text)
            
            # Extract images and perform OCR
            image_list = page.get_images(full=True)
            for img_index, img in enumerate(image_list):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                
                # Perform OCR on the image
                ocr_text = extract_text_from_image(image_bytes)
                if ocr_text.strip():
                    text_content.append(f"\n[Page {page_num+1} Image {img_index+1} Text]:\n{ocr_text}\n")
                    
    except Exception as e:
        return f"Error reading PDF: {e}"
    
    return "\n".join(text_content)

def extract_from_docx(file_bytes: bytes) -> str:
    try:
        doc = Document(io.BytesIO(file_bytes))
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        return "\n".join(full_text)
    except Exception as e:
        return f"Error reading DOCX: {e}"

def extract_from_ppt(file_bytes: bytes) -> str:
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        text_content = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)
                
                # Image extraction from PPT
                if shape.shape_type == 13: # PICTURE
                    try:
                        image_bytes = shape.image.blob
                        ocr_text = extract_text_from_image(image_bytes)
                        if ocr_text.strip():
                            text_content.append(f"[Image Text]: {ocr_text}")
                    except:
                        pass
                        
        return "\n".join(text_content)
    except Exception as e:
        return f"Error reading PPT: {e}"

def extract_from_txt(file_bytes: bytes) -> str:
    try:
        return file_bytes.decode('utf-8')
    except:
        return file_bytes.decode('latin-1')

def extract_content(file_bytes: bytes, filename: str) -> str:
    ext = filename.split('.')[-1].lower()
    
    if ext == 'pdf':
        return extract_from_pdf(file_bytes)
    elif ext in ['docx', 'doc']:
        return extract_from_docx(file_bytes)
    elif ext in ['pptx', 'ppt']:
        return extract_from_ppt(file_bytes)
    elif ext == 'txt':
        return extract_from_txt(file_bytes)
    elif ext in ['jpg', 'jpeg', 'png', 'bmp', 'tiff']:
        return extract_text_from_image(file_bytes)
    else:
        return "Unsupported file format."
